import pandas as pd
import os
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

#generate the plots
def pps_plots():

    x =(0,1,2,3,4,5)
    y =(10,0,250,600,66,85)

    try:
        plt.plot(x, y, "o")
        plt.plot([70, 70], [100, 250], 'k-', lw=2)
        plt.plot([70, 90], [90, 200], 'k-')
    except Exception as e:
        print(e)

    plt.show()
    plt.savefig('graph.jpg')
    img = 'graph.jpg'
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    pic = slide.shapes.add_picture(img, Inches(1), Inches(1), height=Inches(1))



# generate the reports
def report():
    try:
        print('report generated')
    except Exception as e:
        print(e)


